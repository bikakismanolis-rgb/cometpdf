from __future__ import annotations

import argparse
import email
import html
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import textwrap
import tkinter as tk
import zipfile
from io import BytesIO
from pathlib import Path
from tkinter import filedialog, messagebox, ttk
from xml.etree import ElementTree as ET
from email import policy

from PIL import Image, ImageSequence
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape, portrait
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

try:
    from tkinterdnd2 import DND_FILES, TkinterDnD
except Exception:
    DND_FILES = None
    TkinterDnD = None


APP_NAME = "CometPDF"
APP_VERSION = "1.1.1"
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tif", ".tiff", ".webp", ".heic", ".heif"}
TEXT_EXTS = {".txt", ".csv", ".log", ".md", ".html", ".htm", ".xml", ".json"}
WORD_EXTS = {".doc", ".docx", ".rtf", ".odt"}
EXCEL_EXTS = {".xls", ".xlsx", ".xlsm", ".ods", ".csv"}
POWERPOINT_EXTS = {".ppt", ".pptx", ".odp"}
EMAIL_EXTS = {".eml", ".msg"}
SUPPORTED_EXTS = IMAGE_EXTS | TEXT_EXTS | WORD_EXTS | EXCEL_EXTS | POWERPOINT_EXTS | EMAIL_EXTS | {".pdf"}
CONVERSION_MODES = {
    "to_pdf": "Files to PDF",
    "pdf_png": "PDF to PNG images",
    "pdf_jpg": "PDF to JPG images",
    "pdf_txt": "PDF to TXT",
    "pdf_pptx": "PDF to PowerPoint",
    "pdf_docx": "PDF to Word",
}
SETTINGS_FILE = Path(__file__).resolve().parent / "cometpdf_settings.json"


def hidden_subprocess_options() -> dict:
    options: dict = {}
    if os.name == "nt":
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = 0
        options["startupinfo"] = startupinfo
        options["creationflags"] = subprocess.CREATE_NO_WINDOW
    return options


def run_hidden(command: list[str], **kwargs) -> subprocess.CompletedProcess:
    return subprocess.run(command, **hidden_subprocess_options(), **kwargs)


def output_folder() -> Path:
    folder = Path(__file__).resolve().parent / "converted_pdfs"
    folder.mkdir(exist_ok=True)
    return folder


def load_settings() -> dict:
    defaults = {
        "output_dir": str(output_folder()),
        "auto_open": True,
        "open_folder_when_done": False,
        "combine_outputs": False,
        "mode": "to_pdf",
    }
    if not SETTINGS_FILE.exists():
        return defaults
    try:
        loaded = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
        if isinstance(loaded, dict):
            defaults.update({key: loaded[key] for key in defaults if key in loaded})
    except Exception:
        pass
    return defaults


def save_settings(settings: dict) -> None:
    SETTINGS_FILE.write_text(json.dumps(settings, indent=2, ensure_ascii=False), encoding="utf-8")


def default_output_path(input_path: Path) -> Path:
    return input_path.parent / f"{input_path.stem}.pdf"


def mode_output_suffix(mode: str) -> str:
    return {
        "to_pdf": ".pdf",
        "pdf_txt": ".txt",
        "pdf_pptx": ".pptx",
        "pdf_docx": ".docx",
        "pdf_png": ".png",
        "pdf_jpg": ".jpg",
    }.get(mode, ".pdf")


def unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    base = path.with_suffix("")
    suffix = path.suffix
    index = 2
    while True:
        candidate = Path(f"{base}_{index}{suffix}")
        if not candidate.exists():
            return candidate
        index += 1


def merge_pdf_files(pdf_paths: list[Path], output_path: Path) -> Path:
    from pypdf import PdfReader, PdfWriter

    writer = PdfWriter()
    for pdf_path in pdf_paths:
        reader = PdfReader(str(pdf_path))
        for page in reader.pages:
            writer.add_page(page)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("wb") as handle:
        writer.write(handle)
    if not output_path.exists() or output_path.stat().st_size == 0:
        raise RuntimeError("Δεν δημιουργήθηκε έγκυρο ενιαίο PDF.")
    return output_path


def require_pdf(input_path: Path) -> None:
    if input_path.suffix.lower() != ".pdf":
        raise RuntimeError("This output mode only accepts PDF input files.")


def pdf_to_images(input_path: Path, output_base: Path, image_format: str) -> list[Path]:
    require_pdf(input_path)
    try:
        import fitz
    except Exception as exc:
        raise RuntimeError("PDF image export needs PyMuPDF.") from exc

    image_format = image_format.lower()
    if image_format not in {"png", "jpg", "jpeg"}:
        raise RuntimeError("Unsupported image export format.")
    ext = ".jpg" if image_format in {"jpg", "jpeg"} else ".png"

    doc = fitz.open(str(input_path))
    try:
        if doc.page_count == 0:
            raise RuntimeError("The PDF has no pages.")

        output_base.parent.mkdir(parents=True, exist_ok=True)
        outputs: list[Path] = []
        zoom = 2.0
        matrix = fitz.Matrix(zoom, zoom)
        for page_index in range(doc.page_count):
            page = doc.load_page(page_index)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            output_path = unique_path(output_base.parent / f"{output_base.stem}_page_{page_index + 1:03d}{ext}")
            pixmap.save(str(output_path))
            outputs.append(output_path)
        return outputs
    finally:
        doc.close()


def pdf_to_text(input_path: Path, output_path: Path) -> list[Path]:
    require_pdf(input_path)
    from pypdf import PdfReader

    reader = PdfReader(str(input_path))
    parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        parts.append(f"--- Page {index} ---\n{text.strip()}\n")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text("\n".join(parts).strip() + "\n", encoding="utf-8")
    return [output_path]


def pdf_to_pptx(input_path: Path, output_path: Path) -> list[Path]:
    require_pdf(input_path)
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise RuntimeError("PDF to PowerPoint needs PyMuPDF and python-pptx.") from exc

    doc = fitz.open(str(input_path))
    try:
        if doc.page_count == 0:
            raise RuntimeError("The PDF has no pages.")

        first_rect = doc.load_page(0).rect
        prs = Presentation()
        prs.slide_width = Inches(first_rect.width / 72)
        prs.slide_height = Inches(first_rect.height / 72)
        blank = prs.slide_layouts[6]

        with tempfile.TemporaryDirectory(prefix="cometpdf_pdf_pptx_") as tmp:
            tmp_dir = Path(tmp)
            zoom = 2.0
            matrix = fitz.Matrix(zoom, zoom)
            for page_index in range(doc.page_count):
                page = doc.load_page(page_index)
                pixmap = page.get_pixmap(matrix=matrix, alpha=False)
                image_path = tmp_dir / f"page_{page_index + 1:03d}.png"
                pixmap.save(str(image_path))
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return [output_path]
    finally:
        doc.close()


def pdf_to_docx(input_path: Path, output_path: Path) -> list[Path]:
    require_pdf(input_path)
    try:
        import fitz
        from docx import Document
        from docx.shared import Inches
    except Exception as exc:
        raise RuntimeError("PDF to Word needs PyMuPDF and python-docx.") from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(input_path))
    try:
        if doc.page_count == 0:
            raise RuntimeError("The PDF has no pages.")

        word = Document()
        section = word.sections[0]
        section.top_margin = Inches(0)
        section.bottom_margin = Inches(0)
        section.left_margin = Inches(0)
        section.right_margin = Inches(0)

        with tempfile.TemporaryDirectory(prefix="cometpdf_pdf_docx_") as tmp:
            tmp_dir = Path(tmp)
            matrix = fitz.Matrix(2.0, 2.0)
            for page_index in range(doc.page_count):
                page = doc.load_page(page_index)
                rect = page.rect
                if page_index == 0:
                    section.page_width = Inches(rect.width / 72)
                    section.page_height = Inches(rect.height / 72)
                elif page_index > 0:
                    word.add_page_break()

                image_path = tmp_dir / f"page_{page_index + 1:03d}.png"
                page.get_pixmap(matrix=matrix, alpha=False).save(str(image_path))
                paragraph = word.add_paragraph()
                paragraph.paragraph_format.space_after = 0
                run = paragraph.add_run()
                run.add_picture(str(image_path), width=section.page_width)

        word.save(str(output_path))
    finally:
        doc.close()
    return [output_path]


def convert_file(input_file: str | Path, output_file: str | Path | None = None, mode: str = "to_pdf") -> list[Path]:
    input_path = Path(input_file).expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"File not found: {input_path}")

    output_suffix = mode_output_suffix(mode)
    if output_file:
        output_path = Path(output_file).resolve()
    else:
        output_path = input_path.parent / f"{input_path.stem}{output_suffix}"
    output_path = unique_path(output_path)

    if mode == "to_pdf":
        return [convert_to_pdf(input_path, output_path)]
    if mode == "pdf_png":
        return pdf_to_images(input_path, output_path.with_suffix(""), "png")
    if mode == "pdf_jpg":
        return pdf_to_images(input_path, output_path.with_suffix(""), "jpg")
    if mode == "pdf_txt":
        return pdf_to_text(input_path, output_path.with_suffix(".txt"))
    if mode == "pdf_pptx":
        return pdf_to_pptx(input_path, output_path.with_suffix(".pptx"))
    if mode == "pdf_docx":
        return pdf_to_docx(input_path, output_path.with_suffix(".docx"))
    raise RuntimeError(f"Unsupported conversion mode: {mode}")


def operation_status(input_path: Path, mode: str) -> str:
    ext = input_path.suffix.lower()
    if mode == "pdf_png":
        return "Exporting PDF pages to PNG images..."
    if mode == "pdf_jpg":
        return "Exporting PDF pages to JPG images..."
    if mode == "pdf_txt":
        return "Extracting selectable PDF text..."
    if mode == "pdf_pptx":
        return "Creating PowerPoint slides from PDF pages..."
    if mode == "pdf_docx":
        return "Creating a visual Word document from PDF pages..."
    if ext in WORD_EXTS:
        return "Starting Word or LibreOffice export..."
    if ext in EXCEL_EXTS:
        return "Starting Excel or LibreOffice export..."
    if ext in POWERPOINT_EXTS:
        return "Starting PowerPoint or LibreOffice export..."
    if ext in EMAIL_EXTS:
        return "Rendering email content locally..."
    if ext in IMAGE_EXTS:
        return "Rendering image to PDF..."
    if ext in TEXT_EXTS:
        return "Rendering text or data file to PDF..."
    if ext == ".pdf":
        return "Preparing PDF output..."
    return "Converting locally..."


def register_font() -> str:
    candidates = [
        Path(r"C:\Windows\Fonts\arial.ttf"),
        Path(r"C:\Windows\Fonts\segoeui.ttf"),
        Path(r"C:\Windows\Fonts\calibri.ttf"),
    ]
    for font_path in candidates:
        if font_path.exists():
            pdfmetrics.registerFont(TTFont("AppFont", str(font_path)))
            return "AppFont"
    return "Helvetica"


def image_to_pdf(input_path: Path, output_path: Path) -> None:
    frames: list[Image.Image] = []
    try:
        with Image.open(input_path) as image:
            for frame in ImageSequence.Iterator(image):
                prepared = frame.convert("RGB")
                frames.append(prepared.copy())
    except Exception as exc:
        if input_path.suffix.lower() in {".heic", ".heif"}:
            raise RuntimeError(
                "Το HEIC/HEIF χρειάζεται έξτρα υποστήριξη εικόνας στα Windows/Python. "
                "Αν το μετατρέψεις πρώτα σε JPG/PNG, το πρόγραμμα θα το κάνει PDF κανονικά."
            ) from exc
        raise

    if not frames:
        raise RuntimeError("Δεν βρέθηκαν εικόνες μέσα στο αρχείο.")

    first, rest = frames[0], frames[1:]
    first.save(output_path, "PDF", resolution=100.0, save_all=True, append_images=rest)


def text_to_pdf(input_path: Path, output_path: Path) -> None:
    font_name = register_font()
    raw = input_path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "cp1253", "cp1252"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = raw.decode("utf-8", errors="replace")

    pdf = canvas.Canvas(str(output_path), pagesize=A4)
    width, height = A4
    margin = 18 * mm
    line_height = 5.2 * mm
    max_chars = 105

    def new_page() -> float:
        pdf.setFont(font_name, 10)
        return height - margin

    y = new_page()
    for original_line in text.splitlines() or [""]:
        wrapped = textwrap.wrap(
            original_line.expandtabs(4),
            width=max_chars,
            replace_whitespace=False,
            drop_whitespace=False,
        ) or [""]
        for line in wrapped:
            if y < margin:
                pdf.showPage()
                y = new_page()
            pdf.drawString(margin, y, line[:max_chars])
            y -= line_height
    pdf.save()


def eml_to_pdf(input_path: Path, output_path: Path) -> None:
    message = email.message_from_bytes(input_path.read_bytes(), policy=policy.default)
    subject = message.get("subject", "")
    sender = message.get("from", "")
    recipients = message.get("to", "")
    date = message.get("date", "")

    body = ""
    html_body = ""
    if message.is_multipart():
        for part in message.walk():
            content_type = part.get_content_type()
            disposition = part.get_content_disposition()
            if disposition == "attachment":
                continue
            if content_type == "text/plain" and not body:
                body = part.get_content()
            elif content_type == "text/html" and not html_body:
                html_body = part.get_content()
    else:
        if message.get_content_type() == "text/html":
            html_body = message.get_content()
        else:
            body = message.get_content()

    if not body and html_body:
        html_text = re.sub(r"(?i)<br\s*/?>", "\n", html_body)
        html_text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "", html_text)
        html_text = re.sub(r"(?s)<[^>]+>", " ", html_text)
        body = html.unescape("\n".join(line.strip() for line in html_text.splitlines() if line.strip()))

    content = "\n".join(
        [
            f"Subject: {subject}",
            f"From: {sender}",
            f"To: {recipients}",
            f"Date: {date}",
            "",
            body or "(No readable email body found.)",
        ]
    )
    text_to_pdf_from_string(content, output_path)


def msg_to_pdf(input_path: Path, output_path: Path) -> None:
    mhtml_script = r"""
Option Explicit
Dim inputPath, mhtPath, outlook, mail
inputPath = WScript.Arguments.Item(0)
mhtPath = WScript.Arguments.Item(1)
On Error Resume Next
Set outlook = CreateObject("Outlook.Application")
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
Set mail = outlook.Session.OpenSharedItem(inputPath)
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
mail.SaveAs mhtPath, 10
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
On Error GoTo 0
"""
    with tempfile.TemporaryDirectory(prefix="cometpdf_msg_mht_") as tmp:
        mht_path = Path(tmp) / "message.mht"
        try:
            run_vbscript(mhtml_script, [str(input_path), str(mht_path)])
            if mht_path.exists() and mht_path.stat().st_size > 0:
                word_to_pdf(mht_path, output_path)
                if output_path.exists() and output_path.stat().st_size > 0:
                    return
        except Exception:
            pass

    vbs_script = r"""
Option Explicit
Dim inputPath, textPath, outlook, mail, fso, outFile, bodyText
inputPath = WScript.Arguments.Item(0)
textPath = WScript.Arguments.Item(1)
On Error Resume Next
Set outlook = CreateObject("Outlook.Application")
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
Set mail = outlook.Session.OpenSharedItem(inputPath)
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
bodyText = mail.Body
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
Set fso = CreateObject("Scripting.FileSystemObject")
Set outFile = fso.CreateTextFile(textPath, True, True)
outFile.WriteLine "Subject: " & mail.Subject
outFile.WriteLine "From: " & mail.SenderName & " <" & mail.SenderEmailAddress & ">"
outFile.WriteLine "To: " & mail.To
outFile.WriteLine "CC: " & mail.CC
outFile.WriteLine "Date: " & CStr(mail.SentOn)
outFile.WriteLine ""
outFile.WriteLine bodyText
outFile.Close
On Error GoTo 0
"""
    with tempfile.TemporaryDirectory(prefix="cometpdf_msg_") as tmp:
        text_path = Path(tmp) / "message.txt"
        run_vbscript(vbs_script, [str(input_path), str(text_path)])
        if not text_path.exists() or text_path.stat().st_size == 0:
            raise RuntimeError("Το Outlook άνοιξε το MSG, αλλά δεν επέστρεψε περιεχόμενο.")
        content = text_path.read_text(encoding="utf-16", errors="replace")
        text_to_pdf_from_string(content, output_path)


def find_soffice() -> Path | None:
    candidates = [
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return Path(candidate)
    return None


def libreoffice_to_pdf(input_path: Path, output_path: Path) -> bool:
    soffice = find_soffice()
    if not soffice:
        return False

    with tempfile.TemporaryDirectory(prefix="cometpdf_lo_") as tmp:
        result = run_hidden(
            [
                str(soffice),
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                tmp,
                str(input_path),
            ],
            capture_output=True,
            text=True,
            errors="replace",
            timeout=120,
        )
        produced = Path(tmp) / f"{input_path.stem}.pdf"
        if result.returncode != 0 or not produced.exists():
            raise RuntimeError(result.stderr.strip() or result.stdout.strip() or "Το LibreOffice δεν μπόρεσε να κάνει μετατροπή.")
        shutil.move(str(produced), output_path)
    return True


def run_powershell(script: str, args: list[str]) -> None:
    with tempfile.TemporaryDirectory(prefix="cometpdf_ps_") as tmp:
        script_path = Path(tmp) / "convert.ps1"
        script_path.write_text(script, encoding="utf-8")
        result = run_hidden(
            [
                "powershell",
                "-WindowStyle",
                "Hidden",
                "-NoProfile",
                "-ExecutionPolicy",
                "Bypass",
                "-File",
                str(script_path),
                *args,
            ],
            capture_output=True,
            text=True,
            errors="replace",
            timeout=180,
        )
        if result.returncode != 0:
            detail = result.stderr.strip() or result.stdout.strip()
            raise RuntimeError(detail or "Η μετατροπή μέσω Microsoft Office απέτυχε.")


def run_vbscript(script: str, args: list[str]) -> None:
    with tempfile.TemporaryDirectory(prefix="cometpdf_vbs_") as tmp:
        script_path = Path(tmp) / "convert.vbs"
        script_path.write_text(script, encoding="utf-16")
        result = run_hidden(
            ["cscript", "//nologo", str(script_path), *args],
            capture_output=True,
            text=True,
            errors="replace",
            timeout=180,
        )
        if result.returncode != 0:
            detail = result.stderr.strip() or result.stdout.strip()
            raise RuntimeError(detail or "Η μετατροπή μέσω Microsoft Office απέτυχε.")


def xlsx_to_pdf_simple(input_path: Path, output_path: Path) -> None:
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter, range_boundaries
    except ImportError as exc:
        raise RuntimeError("Δεν βρέθηκε το openpyxl για εφεδρική μετατροπή Excel.") from exc

    font_name = register_font()
    workbook = load_workbook(input_path, data_only=True, read_only=False)
    pdf = canvas.Canvas(str(output_path))
    created_page = False

    def cell_text(value: object) -> str:
        if value is None:
            return ""
        text = str(value).replace("\r", "\n")
        return text if len(text) <= 220 else text[:217] + "..."

    def excel_color_to_reportlab(color_obj):
        if not color_obj:
            return None
        rgb = getattr(color_obj, "rgb", None)
        if not rgb or not isinstance(rgb, str):
            return None
        if len(rgb) == 8:
            rgb = rgb[2:]
        if len(rgb) != 6 or rgb.upper() in {"000000", "FFFFFF"}:
            return None
        try:
            return colors.HexColor(f"#{rgb}")
        except ValueError:
            return None

    def column_width_points(sheet, col: int) -> float:
        letter = get_column_letter(col)
        width = sheet.column_dimensions[letter].width or 8.43
        return max(26, min(150, (float(width) * 7 + 5) * 0.75))

    def row_height_points(sheet, row: int) -> float:
        height = sheet.row_dimensions[row].height or sheet.sheet_format.defaultRowHeight or 15
        return max(14, min(80, float(height)))

    def print_bounds(sheet):
        area = sheet.print_area
        if area:
            if isinstance(area, (list, tuple)):
                area = area[0]
            area = str(area).split("!")[-1].replace("$", "").replace("'", "")
            if "," in area:
                area = area.split(",")[0]
            return range_boundaries(area)
        dimension = sheet.calculate_dimension()
        if dimension == "A1:A1" and sheet["A1"].value is None:
            return None
        return range_boundaries(dimension)

    def chunks_by_size(items: list[tuple[int, float]], max_size: float) -> list[list[tuple[int, float]]]:
        chunks = []
        current = []
        total = 0.0
        for item, size in items:
            if current and total + size > max_size:
                chunks.append(current)
                current = []
                total = 0.0
            current.append((item, size))
            total += size
        if current:
            chunks.append(current)
        return chunks

    def merged_anchor_map(sheet):
        anchors = {}
        covered = set()
        for merged in sheet.merged_cells.ranges:
            min_col, min_row, max_col, max_row = merged.bounds
            anchors[(min_row, min_col)] = (min_row, min_col, max_row, max_col)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    if (row, col) != (min_row, min_col):
                        covered.add((row, col))
        return anchors, covered

    def draw_text_in_cell(text: str, x: float, y: float, w: float, h: float, cell, scale: float) -> None:
        if not text:
            return
        font_size = max(5, min(14, float(cell.font.sz or 10) * scale))
        pdf.setFont(font_name, font_size)
        pdf.setFillColor(colors.black)
        lines = []
        for part in text.splitlines() or [""]:
            words = part.split(" ")
            line = ""
            for word in words:
                candidate = word if not line else f"{line} {word}"
                if pdf.stringWidth(candidate, font_name, font_size) <= max(4, w - 4):
                    line = candidate
                else:
                    if line:
                        lines.append(line)
                    line = word
            lines.append(line)

        line_height = font_size * 1.15
        max_lines = max(1, int((h - 4) / line_height))
        lines = lines[:max_lines]
        vertical = cell.alignment.vertical or "center"
        if vertical == "top":
            text_y = y + h - font_size - 3
        elif vertical == "bottom":
            text_y = y + 3 + (len(lines) - 1) * line_height
        else:
            text_y = y + (h + (len(lines) - 1) * line_height) / 2 - font_size * 0.35

        horizontal = cell.alignment.horizontal or "general"
        for line in lines:
            text_width = pdf.stringWidth(line, font_name, font_size)
            if horizontal in {"center", "centerContinuous"}:
                text_x = x + max(2, (w - text_width) / 2)
            elif horizontal == "right":
                text_x = x + max(2, w - text_width - 3)
            else:
                text_x = x + 3
            pdf.drawString(text_x, text_y, line)
            text_y -= line_height

    for sheet_index, sheet in enumerate(workbook.worksheets):
        bounds = print_bounds(sheet)
        if not bounds:
            continue
        min_col, min_row, max_col, max_row = bounds

        page_size = landscape(A4) if sheet.page_setup.orientation == "landscape" else portrait(A4)
        page_width, page_height = page_size
        margins = sheet.page_margins
        left = max(7 * mm, float(margins.left or 0.25) * 72)
        right = max(7 * mm, float(margins.right or 0.25) * 72)
        top = max(7 * mm, float(margins.top or 0.25) * 72)
        bottom = max(7 * mm, float(margins.bottom or 0.25) * 72)
        usable_width = page_width - left - right
        usable_height = page_height - top - bottom

        columns = [
            (col, column_width_points(sheet, col))
            for col in range(min_col, max_col + 1)
            if not sheet.column_dimensions[get_column_letter(col)].hidden
        ]
        rows = [
            (row, row_height_points(sheet, row))
            for row in range(min_row, max_row + 1)
            if not sheet.row_dimensions[row].hidden
        ]
        if not columns or not rows:
            continue

        fit_to_width = int(sheet.page_setup.fitToWidth or 0)
        if fit_to_width == 1:
            column_chunks = [columns]
        else:
            column_chunks = chunks_by_size(columns, usable_width)
        row_chunks = chunks_by_size(rows, usable_height)
        anchors, covered = merged_anchor_map(sheet)

        for row_chunk in row_chunks:
            row_lookup = {row: height for row, height in row_chunk}
            for column_chunk in column_chunks:
                col_lookup = {col: width for col, width in column_chunk}
                content_width = sum(col_lookup.values())
                content_height = sum(row_lookup.values())
                scale = min(1.0, usable_width / content_width, usable_height / content_height)
                drawn_width = content_width * scale
                drawn_height = content_height * scale
                start_x = left
                start_y = page_height - top - drawn_height

                if created_page:
                    pdf.showPage()
                pdf.setPageSize(page_size)
                created_page = True

                y = start_y + drawn_height
                for row, row_height in row_chunk:
                    h = row_height * scale
                    y -= h
                    x = start_x
                    for col, col_width in column_chunk:
                        w = col_width * scale
                        if (row, col) in covered:
                            x += w
                            continue

                        cell = sheet.cell(row=row, column=col)
                        merged = anchors.get((row, col))
                        if merged:
                            _, _, merged_max_row, merged_max_col = merged
                            w = sum(col_lookup.get(c, 0) for c in range(col, merged_max_col + 1)) * scale
                            h = sum(row_lookup.get(r, 0) for r in range(row, merged_max_row + 1)) * scale

                        fill = excel_color_to_reportlab(cell.fill.fgColor if cell.fill else None)
                        if fill and getattr(cell.fill, "fill_type", None) == "solid":
                            pdf.setFillColor(fill)
                            pdf.rect(x, y, w, h, stroke=0, fill=1)

                        pdf.setStrokeColor(colors.HexColor("#9ca3af"))
                        pdf.setLineWidth(0.25)
                        pdf.rect(x, y, w, h, stroke=1, fill=0)
                        draw_text_in_cell(cell_text(cell.value), x, y, w, h, cell, scale)
                        x += col_width * scale

    workbook.close()
    if not created_page:
        text_to_pdf_from_string("Το Excel αρχείο δεν είχε εμφανή δεδομένα.", output_path)
        return
    pdf.save()


def text_to_pdf_from_string(text: str, output_path: Path) -> None:
    font_name = register_font()
    pdf = canvas.Canvas(str(output_path), pagesize=A4)
    width, height = A4
    margin = 18 * mm
    y = height - margin
    pdf.setFont(font_name, 10)
    for line in text.splitlines() or [""]:
        if y < margin:
            pdf.showPage()
            y = height - margin
            pdf.setFont(font_name, 10)
        pdf.drawString(margin, y, line[:100])
        y -= 5.2 * mm
    pdf.save()


def pptx_to_pdf_simple(input_path: Path, output_path: Path) -> None:
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    }

    def rels_for(zip_file: zipfile.ZipFile, slide_name: str) -> dict[str, str]:
        rel_path = f"ppt/slides/_rels/{Path(slide_name).name}.rels"
        if rel_path not in zip_file.namelist():
            return {}
        root = ET.fromstring(zip_file.read(rel_path))
        rels = {}
        for rel in root.findall("rel:Relationship", ns):
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target", "")
            if rel_id and target:
                if target.startswith("../"):
                    target = "ppt/" + target[3:]
                elif not target.startswith("ppt/"):
                    target = "ppt/slides/" + target
                rels[rel_id] = target
        return rels

    def emu_to_pt(value: str | None) -> float:
        return (float(value or 0) / 914400.0) * 72.0

    def presentation_size(zip_file: zipfile.ZipFile) -> tuple[float, float]:
        if "ppt/presentation.xml" not in zip_file.namelist():
            return 960.0, 540.0
        root = ET.fromstring(zip_file.read("ppt/presentation.xml"))
        size = root.find("p:sldSz", ns)
        if size is None:
            return 960.0, 540.0
        width = emu_to_pt(size.attrib.get("cx"))
        height = emu_to_pt(size.attrib.get("cy"))
        if width <= 0 or height <= 0:
            return 960.0, 540.0
        return width, height

    def shape_box(shape) -> tuple[float, float, float, float]:
        xfrm = shape.find(".//a:xfrm", ns)
        if xfrm is None:
            return 40, 40, 400, 80
        off = xfrm.find("a:off", ns)
        ext = xfrm.find("a:ext", ns)
        return (
            emu_to_pt(off.attrib.get("x") if off is not None else None),
            emu_to_pt(off.attrib.get("y") if off is not None else None),
            emu_to_pt(ext.attrib.get("cx") if ext is not None else None),
            emu_to_pt(ext.attrib.get("cy") if ext is not None else None),
        )

    def shape_text(shape) -> str:
        paragraphs = []
        for paragraph in shape.findall(".//a:p", ns):
            parts = [node.text or "" for node in paragraph.findall(".//a:t", ns)]
            if parts:
                paragraphs.append("".join(parts))
        return "\n".join(paragraphs).strip()

    def draw_wrapped(pdf: canvas.Canvas, text: str, x: float, y: float, w: float, h: float, scale: float) -> None:
        if not text:
            return
        font_name = register_font()
        font_size = max(7, min(22, 13 * scale))
        pdf.setFont(font_name, font_size)
        pdf.setFillColor(colors.black)
        line_height = font_size * 1.18
        max_chars = max(12, int(w / max(font_size * 0.52, 1)))
        lines = []
        for line in text.splitlines():
            lines.extend(textwrap.wrap(line, width=max_chars) or [""])
        max_lines = max(1, int(h / line_height))
        cursor_y = y + h - font_size
        for line in lines[:max_lines]:
            pdf.drawString(x, cursor_y, line)
            cursor_y -= line_height

    with zipfile.ZipFile(input_path) as pptx:
        slide_names = sorted(
            [name for name in pptx.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")],
            key=lambda value: int("".join(ch for ch in Path(value).stem if ch.isdigit()) or "0"),
        )
        if not slide_names:
            raise RuntimeError("Το PPTX δεν περιέχει slides.")

        slide_width, slide_height = presentation_size(pptx)
        page_width, page_height = slide_width, slide_height
        pdf = canvas.Canvas(str(output_path), pagesize=(page_width, page_height))
        scale = min(page_width / slide_width, page_height / slide_height)
        x_offset = (page_width - slide_width * scale) / 2
        y_offset = (page_height - slide_height * scale) / 2

        for index, slide_name in enumerate(slide_names):
            if index:
                pdf.showPage()
            pdf.setPageSize((page_width, page_height))
            pdf.setFillColor(colors.white)
            pdf.rect(0, 0, page_width, page_height, stroke=0, fill=1)
            root = ET.fromstring(pptx.read(slide_name))
            rels = rels_for(pptx, slide_name)

            for pic in root.findall(".//p:pic", ns):
                x, y, w, h = shape_box(pic)
                blip = pic.find(".//a:blip", ns)
                rel_id = blip.attrib.get(f"{{{ns['r']}}}embed") if blip is not None else None
                image_path = rels.get(rel_id or "")
                if image_path and image_path in pptx.namelist():
                    try:
                        image = Image.open(BytesIO(pptx.read(image_path))).convert("RGB")
                        draw_x = x_offset + x * scale
                        draw_y = page_height - y_offset - (y + h) * scale
                        pdf.drawInlineImage(image, draw_x, draw_y, w * scale, h * scale)
                    except Exception:
                        pass

            for shape in root.findall(".//p:sp", ns):
                text = shape_text(shape)
                if not text:
                    continue
                x, y, w, h = shape_box(shape)
                draw_x = x_offset + x * scale
                draw_y = page_height - y_offset - (y + h) * scale
                draw_wrapped(pdf, text, draw_x, draw_y, w * scale, h * scale, scale)

        pdf.save()


def docx_to_pdf_simple(input_path: Path, output_path: Path) -> None:
    ns = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    }

    def rels_for_doc(zip_file: zipfile.ZipFile) -> dict[str, str]:
        rel_path = "word/_rels/document.xml.rels"
        if rel_path not in zip_file.namelist():
            return {}
        root = ET.fromstring(zip_file.read(rel_path))
        rels = {}
        for rel in root.findall("rel:Relationship", ns):
            rel_id = rel.attrib.get("Id")
            target = rel.attrib.get("Target", "")
            if rel_id and target:
                if target.startswith("../"):
                    target = target[3:]
                elif not target.startswith("word/"):
                    target = "word/" + target
                rels[rel_id] = target
        return rels

    def draw_paragraph(pdf: canvas.Canvas, text: str, x: float, y: float, width: float, font_name: str) -> float:
        if not text.strip():
            return y - 7
        font_size = 10.5
        line_height = 14
        pdf.setFont(font_name, font_size)
        pdf.setFillColor(colors.black)
        max_chars = max(20, int(width / 5.2))
        for line in textwrap.wrap(text, width=max_chars) or [""]:
            pdf.drawString(x, y, line)
            y -= line_height
        return y - 3

    font_name = register_font()
    pdf = canvas.Canvas(str(output_path), pagesize=A4)
    page_width, page_height = A4
    margin = 18 * mm
    y = page_height - margin

    with zipfile.ZipFile(input_path) as docx:
        if "word/document.xml" not in docx.namelist():
            raise RuntimeError("Το DOCX δεν περιέχει έγγραφο Word.")
        rels = rels_for_doc(docx)
        root = ET.fromstring(docx.read("word/document.xml"))

        for child in root.findall(".//w:body/*", ns):
            if child.tag.endswith("}p"):
                text = "".join(node.text or "" for node in child.findall(".//w:t", ns))
                if y < margin + 25:
                    pdf.showPage()
                    y = page_height - margin
                y = draw_paragraph(pdf, text, margin, y, page_width - 2 * margin, font_name)

                for blip in child.findall(".//a:blip", ns):
                    rel_id = blip.attrib.get(f"{{{ns['r']}}}embed")
                    image_path = rels.get(rel_id or "")
                    if image_path and image_path in docx.namelist():
                        try:
                            image = Image.open(BytesIO(docx.read(image_path))).convert("RGB")
                            max_w = page_width - 2 * margin
                            max_h = page_height * 0.45
                            scale = min(max_w / image.width, max_h / image.height, 1.0)
                            draw_w = image.width * scale
                            draw_h = image.height * scale
                            if y - draw_h < margin:
                                pdf.showPage()
                                y = page_height - margin
                            pdf.drawInlineImage(image, margin, y - draw_h, draw_w, draw_h)
                            y -= draw_h + 10
                        except Exception:
                            pass

            elif child.tag.endswith("}tbl"):
                rows = []
                for row in child.findall(".//w:tr", ns):
                    cells = []
                    for cell in row.findall("w:tc", ns):
                        cell_text = " ".join("".join(node.text or "" for node in p.findall(".//w:t", ns)) for p in cell.findall("w:p", ns))
                        cells.append(cell_text.strip())
                    if cells:
                        rows.append(cells)
                for row in rows:
                    if y < margin + 18:
                        pdf.showPage()
                        y = page_height - margin
                    y = draw_paragraph(pdf, " | ".join(row), margin, y, page_width - 2 * margin, font_name)

    pdf.save()


def word_to_pdf(input_path: Path, output_path: Path) -> None:
    vbs_script = r"""
Option Explicit
Dim inputPath, outputPath, word, doc
inputPath = WScript.Arguments.Item(0)
outputPath = WScript.Arguments.Item(1)
Set word = CreateObject("Word.Application")
word.Visible = False
Set doc = word.Documents.Open(inputPath, False, True)
doc.ExportAsFixedFormat outputPath, 17
doc.Close False
word.Quit
"""
    try:
        run_vbscript(vbs_script, [str(input_path), str(output_path)])
        if output_path.exists() and output_path.stat().st_size > 0:
            return
    except Exception:
        pass

    script = r"""
param([string]$InputPath, [string]$OutputPath)
$ErrorActionPreference = 'Stop'
$word = $null
$doc = $null
try {
  $word = New-Object -ComObject Word.Application
  if ($null -eq $word) {
    throw "Το Word COM δεν ξεκίνησε."
  }
  try { $word.Visible = $false } catch { }
  $doc = $word.Documents.Open($InputPath, $false, $true)
  $doc.ExportAsFixedFormat($OutputPath, 17)
  if (-not (Test-Path -LiteralPath $OutputPath)) {
    throw "Το Word δεν δημιούργησε PDF."
  }
} finally {
  if ($doc -ne $null) { $doc.Close($false) }
  if ($word -ne $null) {
    try { $word.Quit() } catch { }
  }
}
"""
    try:
        run_powershell(script, [str(input_path), str(output_path)])
    except Exception as exc:
        if input_path.suffix.lower() == ".docx":
            docx_to_pdf_simple(input_path, output_path)
            return
        raise RuntimeError(f"Το Microsoft Word δεν μπόρεσε να κάνει εξαγωγή σε PDF: {exc}") from exc


def excel_to_pdf(input_path: Path, output_path: Path) -> None:
    vbs_script = r"""
Option Explicit
Dim inputPath, outputPath, excel, book
inputPath = WScript.Arguments.Item(0)
outputPath = WScript.Arguments.Item(1)
On Error Resume Next
Set excel = CreateObject("Excel.Application")
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
excel.Visible = False
excel.DisplayAlerts = False
Set book = excel.Workbooks.Open(inputPath, 3, True)
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  excel.Quit
  WScript.Quit 1
End If
book.ExportAsFixedFormat 0, outputPath
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  book.Close False
  excel.Quit
  WScript.Quit 1
End If
book.Close False
excel.Quit
On Error GoTo 0
"""
    try:
        run_vbscript(vbs_script, [str(input_path), str(output_path)])
        if output_path.exists() and output_path.stat().st_size > 0:
            return
    except Exception:
        pass

    script = r"""
param([string]$InputPath, [string]$OutputPath)
$ErrorActionPreference = 'Stop'
$excel = $null
$book = $null
try {
  $excel = New-Object -ComObject Excel.Application
  $excel.Visible = $false
  $excel.DisplayAlerts = $false
  $book = $excel.Workbooks.Open($InputPath, 3, $true)
  $book.ExportAsFixedFormat(0, $OutputPath)
  if (-not (Test-Path -LiteralPath $OutputPath)) {
    throw "Το Excel δεν δημιούργησε PDF."
  }
} finally {
  if ($book -ne $null) { $book.Close($false) }
  if ($excel -ne $null) {
    try { $excel.Quit() } catch { }
  }
}
"""
    try:
        run_powershell(script, [str(input_path), str(output_path)])
    except Exception as exc:
        if input_path.suffix.lower() in {".xlsx", ".xlsm"}:
            xlsx_to_pdf_simple(input_path, output_path)
            return
        raise RuntimeError(f"Το Microsoft Excel δεν μπόρεσε να κάνει εξαγωγή σε PDF: {exc}") from exc


def powerpoint_to_pdf(input_path: Path, output_path: Path) -> None:
    vbs_script = r"""
Option Explicit
Dim inputPath, outputPath, ppt, presentation
inputPath = WScript.Arguments.Item(0)
outputPath = WScript.Arguments.Item(1)
On Error Resume Next
Set ppt = CreateObject("PowerPoint.Application")
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  WScript.Quit 1
End If
Set presentation = ppt.Presentations.Open(inputPath, True, True, False)
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  ppt.Quit
  WScript.Quit 1
End If
presentation.SaveAs outputPath, 32
If Err.Number <> 0 Then
  WScript.Echo Err.Description
  presentation.Close
  ppt.Quit
  WScript.Quit 1
End If
presentation.Close
ppt.Quit
On Error GoTo 0
"""
    try:
        run_vbscript(vbs_script, [str(input_path), str(output_path)])
        if output_path.exists() and output_path.stat().st_size > 0:
            return
    except Exception:
        pass

    script = r"""
param([string]$InputPath, [string]$OutputPath)
$ErrorActionPreference = 'Stop'
$ppt = $null
$presentation = $null
try {
  $ppt = New-Object -ComObject PowerPoint.Application
  if ($null -eq $ppt) {
    throw "Το PowerPoint COM δεν ξεκίνησε."
  }
  $presentation = $ppt.Presentations.Open($InputPath, $true, $true, $false)
  if ($null -eq $presentation) {
    throw "Το PowerPoint δεν μπόρεσε να ανοίξει την παρουσίαση."
  }
  $presentation.SaveAs($OutputPath, 32)
  if (-not (Test-Path -LiteralPath $OutputPath)) {
    throw "Το PowerPoint δεν δημιούργησε PDF."
  }
} finally {
  if ($presentation -ne $null) { $presentation.Close() }
  if ($ppt -ne $null) {
    try { $ppt.Quit() } catch { }
  }
}
"""
    try:
        run_powershell(script, [str(input_path), str(output_path)])
    except Exception as exc:
        if input_path.suffix.lower() == ".pptx":
            pptx_to_pdf_simple(input_path, output_path)
            return
        raise RuntimeError(f"Το Microsoft PowerPoint δεν μπόρεσε να κάνει εξαγωγή σε PDF: {exc}") from exc


def office_to_pdf(input_path: Path, output_path: Path) -> None:
    if libreoffice_to_pdf(input_path, output_path):
        return

    ext = input_path.suffix.lower()
    if ext in WORD_EXTS:
        word_to_pdf(input_path, output_path)
    elif ext in EXCEL_EXTS:
        excel_to_pdf(input_path, output_path)
    elif ext in POWERPOINT_EXTS:
        powerpoint_to_pdf(input_path, output_path)
    else:
        raise RuntimeError("Δεν υπάρχει διαθέσιμος μετατροπέας για αυτόν τον τύπο αρχείου.")


def convert_to_pdf(input_file: str | Path, output_file: str | Path | None = None) -> Path:
    input_path = Path(input_file).expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"Δεν βρέθηκε το αρχείο: {input_path}")

    ext = input_path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise RuntimeError(f"Δεν υποστηρίζεται ακόμη ο τύπος αρχείου: {ext or '(χωρίς κατάληξη)'}")

    output_path = Path(output_file).resolve() if output_file else unique_path(default_output_path(input_path))
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if ext == ".pdf":
        if input_path != output_path:
            shutil.copy2(input_path, output_path)
    elif ext in IMAGE_EXTS:
        image_to_pdf(input_path, output_path)
    elif ext in TEXT_EXTS:
        text_to_pdf(input_path, output_path)
    elif ext == ".eml":
        eml_to_pdf(input_path, output_path)
    elif ext == ".msg":
        msg_to_pdf(input_path, output_path)
    else:
        office_to_pdf(input_path, output_path)

    if not output_path.exists() or output_path.stat().st_size == 0:
        raise RuntimeError("Η μετατροπή δεν δημιούργησε έγκυρο PDF.")
    return output_path


BaseTk = TkinterDnD.Tk if TkinterDnD else tk.Tk


class CometPDFApp(BaseTk):
    def __init__(self) -> None:
        super().__init__()
        self.title(APP_NAME)
        self.geometry("920x740")
        self.minsize(580, 460)
        self.configure(bg="#f5f7fb")

        self.settings = load_settings()
        self.files: list[Path] = []
        self.outputs: list[Path] = []
        self.is_running = False

        saved_mode = self.settings.get("mode", "to_pdf")
        if saved_mode not in CONVERSION_MODES:
            saved_mode = "to_pdf"

        self.output_dir = tk.StringVar(value=self.settings["output_dir"])
        self.mode = tk.StringVar(value=saved_mode)
        self.status = tk.StringVar(value="Add one or more files and choose a conversion mode.")
        self.count_label = tk.StringVar(value="0 files")
        self.auto_open = tk.BooleanVar(value=bool(self.settings["auto_open"]))
        self.open_folder_when_done = tk.BooleanVar(value=bool(self.settings["open_folder_when_done"]))
        self.combine_outputs = tk.BooleanVar(value=bool(self.settings["combine_outputs"]))
        self.progress_value = tk.DoubleVar(value=0)

        self._build_ui()

    def _build_ui(self) -> None:
        style = ttk.Style(self)
        style.theme_use("clam")
        style.configure("App.TFrame", background="#f5f7fb")
        style.configure("Panel.TFrame", background="#ffffff", relief="solid", borderwidth=1)
        style.configure("Title.TLabel", background="#f5f7fb", font=("Segoe UI", 22, "bold"))
        style.configure("Subtle.TLabel", background="#f5f7fb", foreground="#475569", font=("Segoe UI", 10))
        style.configure("Panel.TLabel", background="#ffffff", font=("Segoe UI", 10))
        style.configure("Status.TLabel", background="#eef2f7", foreground="#0f172a", font=("Segoe UI", 10), padding=12)
        style.configure("TButton", font=("Segoe UI", 10), padding=(12, 7))
        style.configure("Accent.TButton", font=("Segoe UI", 12, "bold"), padding=(18, 10))

        shell = ttk.Frame(self, style="App.TFrame")
        shell.pack(fill="both", expand=True)
        self.scroll_canvas = tk.Canvas(shell, bg="#f5f7fb", highlightthickness=0)
        page_scrollbar = ttk.Scrollbar(shell, orient="vertical", command=self.scroll_canvas.yview)
        self.scroll_canvas.configure(yscrollcommand=page_scrollbar.set)
        self.scroll_canvas.pack(side="left", fill="both", expand=True)
        page_scrollbar.pack(side="right", fill="y")

        root = ttk.Frame(self.scroll_canvas, padding=24, style="App.TFrame")
        self.scroll_window = self.scroll_canvas.create_window((0, 0), window=root, anchor="nw")
        root.bind("<Configure>", self.on_scroll_content_configure)
        self.scroll_canvas.bind("<Configure>", self.on_scroll_canvas_configure)
        self.bind_all("<MouseWheel>", self.on_mousewheel)
        self.bind_all("<Button-4>", self.on_mousewheel)
        self.bind_all("<Button-5>", self.on_mousewheel)

        header = ttk.Frame(root, style="App.TFrame")
        header.pack(fill="x")
        ttk.Label(header, text="CometPDF", style="Title.TLabel").pack(anchor="w")
        ttk.Label(
            header,
            text="Convert files to PDF, or export PDFs to images, text, Word, and PowerPoint.",
            style="Subtle.TLabel",
        ).pack(anchor="w", pady=(2, 18))

        picker = ttk.Frame(root, padding=18, style="Panel.TFrame")
        picker.pack(fill="x", pady=(0, 14))
        ttk.Label(picker, text="Add files", style="Panel.TLabel", font=("Segoe UI", 13, "bold")).pack(anchor="w")
        ttk.Label(
            picker,
            text="Drag files here or select several files at once. Results are saved in the output folder.",
            style="Panel.TLabel",
        ).pack(anchor="w", pady=(3, 12))

        picker_buttons = ttk.Frame(picker, style="Panel.TFrame")
        picker_buttons.pack(fill="x")
        ttk.Button(picker_buttons, text="Add files...", command=self.choose_files).pack(side="left")
        ttk.Button(picker_buttons, text="Remove selected", command=self.remove_selected).pack(side="left", padx=(8, 0))
        ttk.Button(picker_buttons, text="Clear", command=self.clear_files).pack(side="left", padx=(8, 0))
        ttk.Label(picker_buttons, textvariable=self.count_label, style="Panel.TLabel").pack(side="right")

        list_frame = ttk.Frame(root, style="App.TFrame")
        list_frame.pack(fill="both", expand=True, pady=(0, 14))
        self.file_list = tk.Listbox(
            list_frame,
            selectmode=tk.EXTENDED,
            font=("Segoe UI", 10),
            activestyle="none",
            bg="#ffffff",
            fg="#0f172a",
            highlightthickness=1,
            highlightbackground="#cbd5e1",
            relief="flat",
            height=9,
        )
        self.file_list.pack(side="left", fill="both", expand=True)
        scrollbar = ttk.Scrollbar(list_frame, orient="vertical", command=self.file_list.yview)
        scrollbar.pack(side="right", fill="y")
        self.file_list.configure(yscrollcommand=scrollbar.set)
        self.enable_drag_drop(root, picker, self.file_list)

        output = ttk.Frame(root, padding=14, style="Panel.TFrame")
        output.pack(fill="x", pady=(0, 14))
        ttk.Label(output, text="Output folder", style="Panel.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Entry(output, textvariable=self.output_dir).grid(row=1, column=0, sticky="ew", pady=(6, 0), padx=(0, 8))
        ttk.Button(output, text="Change...", command=self.choose_output_dir).grid(row=1, column=1, sticky="e", pady=(6, 0))
        output.columnconfigure(0, weight=1)

        mode_box = ttk.Frame(root, padding=14, style="Panel.TFrame")
        mode_box.pack(fill="x", pady=(0, 14))
        ttk.Label(mode_box, text="Conversion mode", style="Panel.TLabel").grid(row=0, column=0, sticky="w")
        self.mode_combo = ttk.Combobox(mode_box, state="readonly", values=list(CONVERSION_MODES.values()), width=34)
        self.mode_combo.set(CONVERSION_MODES[self.mode.get()])
        self.mode_combo.grid(row=1, column=0, sticky="w", pady=(6, 0))
        self.mode_combo.bind("<<ComboboxSelected>>", self.on_mode_changed)
        ttk.Label(
            mode_box,
            text="Use Files to PDF for normal conversion. Use PDF modes only when the input is a PDF.",
            style="Panel.TLabel",
        ).grid(row=2, column=0, sticky="w", pady=(8, 0))

        options = ttk.Frame(root, style="App.TFrame")
        options.pack(fill="x", pady=(0, 12))
        ttk.Checkbutton(options, text="Open result after single-file conversion", variable=self.auto_open, command=self.save_current_settings).pack(anchor="w")
        ttk.Checkbutton(options, text="Open output folder when done", variable=self.open_folder_when_done, command=self.save_current_settings).pack(anchor="w", pady=(3, 0))
        self.combine_check = ttk.Checkbutton(options, text="Merge all PDF outputs into one PDF", variable=self.combine_outputs, command=self.save_current_settings)
        self.combine_check.pack(anchor="w", pady=(3, 0))

        actions = ttk.Frame(root, style="App.TFrame")
        actions.pack(fill="x", pady=(0, 12))
        self.convert_button = ttk.Button(actions, text="Convert", style="Accent.TButton", command=self.convert_all)
        self.convert_button.grid(row=0, column=0, sticky="w", padx=(0, 8), pady=(0, 8))
        ttk.Button(actions, text="Open last result", command=self.open_last_pdf).grid(row=0, column=1, sticky="w", padx=(0, 8), pady=(0, 8))
        ttk.Button(actions, text="Open folder", command=self.open_output_folder).grid(row=0, column=2, sticky="w", padx=(0, 8), pady=(0, 8))
        ttk.Button(actions, text="About", command=self.show_about).grid(row=0, column=3, sticky="w", pady=(0, 8))

        self.progress = ttk.Progressbar(root, variable=self.progress_value, maximum=100)
        self.progress.pack(fill="x", pady=(0, 10))
        ttk.Label(root, textvariable=self.status, style="Status.TLabel", wraplength=820).pack(fill="x")

        supported = "Supports: Office, PDF, images, email, TXT/CSV/HTML/XML/JSON. PDF export: PNG, JPG, TXT, PPTX, DOCX."
        ttk.Label(root, text=supported, style="Subtle.TLabel").pack(anchor="w", pady=(10, 0))
        self.update_mode_state()

    def on_scroll_content_configure(self, _event=None) -> None:
        self.scroll_canvas.configure(scrollregion=self.scroll_canvas.bbox("all"))

    def on_scroll_canvas_configure(self, event) -> None:
        self.scroll_canvas.itemconfigure(self.scroll_window, width=event.width)

    def on_mousewheel(self, event) -> None:
        if not hasattr(self, "scroll_canvas"):
            return
        if getattr(event, "num", None) == 4:
            self.scroll_canvas.yview_scroll(-3, "units")
        elif getattr(event, "num", None) == 5:
            self.scroll_canvas.yview_scroll(3, "units")
        else:
            delta = int(-1 * (event.delta / 120)) if event.delta else 0
            if delta:
                self.scroll_canvas.yview_scroll(delta * 3, "units")

    def enable_drag_drop(self, *widgets) -> None:
        if DND_FILES is None:
            return
        for widget in widgets:
            try:
                widget.drop_target_register(DND_FILES)
                widget.dnd_bind("<<Drop>>", self.on_drop_files)
            except Exception:
                pass

    def on_drop_files(self, event) -> None:
        try:
            paths = self.tk.splitlist(event.data)
        except Exception:
            paths = [event.data]
        self.add_files(paths)

    def selected_mode_key(self) -> str:
        selected = self.mode_combo.get() if hasattr(self, "mode_combo") else CONVERSION_MODES[self.mode.get()]
        for key, label in CONVERSION_MODES.items():
            if label == selected:
                return key
        return "to_pdf"

    def on_mode_changed(self, _event=None) -> None:
        self.mode.set(self.selected_mode_key())
        self.update_mode_state()
        self.save_current_settings()

    def update_mode_state(self) -> None:
        mode = self.mode.get()
        if hasattr(self, "combine_check"):
            self.combine_check.configure(state="normal" if mode == "to_pdf" else "disabled")
        if hasattr(self, "convert_button"):
            self.convert_button.configure(text="Convert to PDF" if mode == "to_pdf" else "Export PDF")
        if mode == "to_pdf":
            self.status.set("Ready for file-to-PDF conversion.")
        else:
            self.status.set(f"Ready for {CONVERSION_MODES[mode]}. Add PDF files only.")

    def save_current_settings(self) -> None:
        self.settings["output_dir"] = self.output_dir.get()
        self.settings["auto_open"] = bool(self.auto_open.get())
        self.settings["open_folder_when_done"] = bool(self.open_folder_when_done.get())
        self.settings["combine_outputs"] = bool(self.combine_outputs.get())
        self.settings["mode"] = self.mode.get()
        save_settings(self.settings)

    def choose_files(self) -> None:
        paths = filedialog.askopenfilenames(
            title="Add files",
            filetypes=[
                ("Supported files", " ".join(f"*{ext}" for ext in sorted(SUPPORTED_EXTS))),
                ("All files", "*.*"),
            ],
        )
        self.add_files(paths)

    def add_files(self, paths) -> None:
        existing = {path.resolve() for path in self.files}
        for item in paths:
            path = Path(item).expanduser().resolve()
            if path.exists() and path.resolve() not in existing:
                self.files.append(path)
                existing.add(path.resolve())
                self.file_list.insert(tk.END, str(path))
        self.refresh_count()

    def remove_selected(self) -> None:
        selected = list(self.file_list.curselection())
        for index in reversed(selected):
            self.file_list.delete(index)
            del self.files[index]
        self.refresh_count()

    def clear_files(self) -> None:
        if self.is_running:
            return
        self.files.clear()
        self.file_list.delete(0, tk.END)
        self.outputs.clear()
        self.progress_value.set(0)
        self.refresh_count()
        self.status.set("The file list was cleared.")

    def refresh_count(self) -> None:
        count = len(self.files)
        self.count_label.set(f"{count} file" if count == 1 else f"{count} files")
        if count:
            self.status.set("Ready to convert.")

    def choose_output_dir(self) -> None:
        folder = filedialog.askdirectory(title="Output folder", initialdir=self.output_dir.get() or str(output_folder()))
        if folder:
            self.output_dir.set(folder)
            self.save_current_settings()

    def convert_all(self) -> None:
        if self.is_running:
            return
        if not self.files:
            messagebox.showwarning(APP_NAME, "Add one or more files first.")
            return

        mode = self.mode.get()
        if mode != "to_pdf" and any(path.suffix.lower() != ".pdf" for path in self.files):
            messagebox.showwarning(APP_NAME, "PDF export modes only accept PDF input files.")
            return

        out_dir = Path(self.output_dir.get()).expanduser()
        try:
            out_dir.mkdir(parents=True, exist_ok=True)
        except Exception as exc:
            messagebox.showerror(APP_NAME, f"Could not open the output folder:\n{exc}")
            return

        self.save_current_settings()
        self.is_running = True
        self.outputs.clear()
        self.convert_button.configure(state="disabled")
        self.progress_value.set(0)
        self.status.set("Starting conversion...")

        worker = threading.Thread(
            target=self._convert_worker,
            args=(list(self.files), out_dir, bool(self.auto_open.get()), bool(self.combine_outputs.get()), mode),
            daemon=True,
        )
        worker.start()

    def _convert_worker(self, files: list[Path], out_dir: Path, auto_open: bool, combine_outputs: bool, mode: str) -> None:
        errors: list[str] = []
        converted: list[Path] = []
        total = len(files)
        for index, input_path in enumerate(files, start=1):
            self.after(0, self._set_status, f"{index}/{total}: {input_path.name} - {operation_status(input_path, mode)}")
            output_path = unique_path(out_dir / f"{input_path.stem}{mode_output_suffix(mode)}")
            try:
                results = convert_file(input_path, output_path, mode)
                converted.extend(results)
                self.outputs.extend(results)
                if auto_open and total == 1:
                    if len(results) == 1:
                        os.startfile(results[0])
                    else:
                        os.startfile(out_dir)
            except Exception as exc:
                errors.append(f"{input_path.name}: {self.human_error(exc)}")
            self.after(0, self.progress_value.set, (index / total) * 100)

        if mode == "to_pdf" and combine_outputs and len(converted) > 1:
            try:
                self.after(0, self._set_status, "Merging PDFs into one file...")
                merged = merge_pdf_files(converted, unique_path(out_dir / "combined.pdf"))
                self.outputs.append(merged)
                if auto_open:
                    os.startfile(merged)
            except Exception as exc:
                errors.append(f"Merge PDFs: {self.human_error(exc)}")

        self.after(0, self._finish_conversion, errors)

    def _set_status(self, text: str) -> None:
        self.status.set(text)

    def _finish_conversion(self, errors: list[str]) -> None:
        self.is_running = False
        self.convert_button.configure(state="normal")
        self.update_mode_state()
        success_count = len(self.outputs)
        if errors:
            self.status.set(f"Completed with {success_count} results and {len(errors)} errors.")
            messagebox.showwarning(APP_NAME, "Some files were not converted:\n\n" + "\n".join(errors[:8]))
        else:
            self.status.set(f"Done. Created {success_count} result file(s).")
        if self.open_folder_when_done.get() and success_count:
            self.open_output_folder()

    def human_error(self, exc: Exception) -> str:
        text = str(exc).strip() or exc.__class__.__name__
        replacements = {
            "COM": "Microsoft Office",
            "VBScript": "Office automation",
            "Invalid procedure call or argument": "The source application could not export this file type.",
        }
        for old, new in replacements.items():
            text = text.replace(old, new)
        return text

    def open_last_pdf(self) -> None:
        if self.outputs and self.outputs[-1].exists():
            os.startfile(self.outputs[-1])
        else:
            messagebox.showinfo(APP_NAME, "There is no result file yet.")

    def open_output_folder(self) -> None:
        folder = Path(self.output_dir.get()).expanduser()
        folder.mkdir(parents=True, exist_ok=True)
        os.startfile(folder)

    def show_about(self) -> None:
        messagebox.showinfo(
            APP_NAME,
            f"{APP_NAME} {APP_VERSION}\n\n"
            "Convert files to PDF and export PDFs to images, text, Word, and PowerPoint.\n\n"
            "Office conversion uses Microsoft Office, Outlook, or LibreOffice when available. "
            "PDF export uses local processing on this computer.",
        )

def main() -> int:
    parser = argparse.ArgumentParser(description="Simple file to PDF converter.")
    parser.add_argument("--convert", help="Input file to convert without opening the GUI.")
    parser.add_argument("--output", help="Output PDF path for --convert.")
    parser.add_argument(
        "--mode",
        choices=sorted(CONVERSION_MODES),
        default="to_pdf",
        help="Conversion mode for --convert.",
    )
    parser.add_argument("--open", action="store_true", help="Open the PDF after converting.")
    args = parser.parse_args()

    if args.convert:
        try:
            outputs = convert_file(args.convert, args.output, args.mode)
            if args.open:
                if len(outputs) == 1:
                    os.startfile(outputs[0])
                elif outputs:
                    os.startfile(outputs[0].parent)
            print("\n".join(str(path) for path in outputs))
            return 0
        except Exception as exc:
            print(exc, file=sys.stderr)
            return 1

    app = CometPDFApp()
    app.mainloop()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
